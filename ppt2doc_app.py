import os
import shutil
import tempfile
import pandas as pd
import streamlit as st
import openai
from pptx import Presentation
from PIL import Image
from pdf2image import convert_from_path
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches

# Add a text input for the OpenAI API key
openai_api_key = st.text_input("Enter your OpenAI API key:", type="password")

# Check if the API key is provided
if openai_api_key:
    openai.api_key = openai_api_key

    # Ensure poppler is in the PATH
    os.environ["PATH"] += os.pathsep + r'C:\Users\User\Dropbox\PC\Documents\work\Peripeteia\poppler-24.02.0\Library\bin'

    # Function to ensure the existence of an image directory and clear it if it exists
    def ensure_image_directory():
        if os.path.exists("extracted_images"):
            try:
                shutil.rmtree("extracted_images")
            except PermissionError:
                pass
        os.makedirs("extracted_images", exist_ok=True)

    # Function to sanitize text by removing non-printable characters
    def sanitize_text(text):
        return ''.join(c for c in text if c.isprintable())

    # Function to extract and crop images from PDF
    def extract_images_from_pdf(pdf_path, crop_box=None):
        images = convert_from_path(pdf_path)
        image_paths = []
        for i, image in enumerate(images):
            if crop_box:
                image = image.crop(crop_box)
            image_path = f"extracted_images/slide_{i + 1}.png"
            image.save(image_path, 'PNG')
            image_paths.append(image_path)
        return image_paths

    # Function to extract text from a PowerPoint file
    def extract_text_from_pptx(file_path):
        prs = Presentation(file_path)
        slides_content = []
        for i, slide in enumerate(prs.slides):
            slide_title = ""
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if not slide_title and shape.text_frame:
                        slide_title = shape.text_frame.text
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            slide_text += paragraph.text + "\n"
            slides_content.append({
                "title": sanitize_text(slide_title.strip()),
                "text": sanitize_text(slide_text.strip())
            })
        return slides_content

    def generate_section_text(slide_content, section, presentation_context, example_report=None):
        combined_text = "\n\n".join([sanitize_text(slide["text"]) for slide in slide_content])
        combined_titles = " ".join([sanitize_text(slide["title"]) for slide in slide_content])
        
        prompt = (
            f"Using the context of '{presentation_context}', write a technical description that summarizes the content of the following slides. "
            f"The description should be clear, concise, and fit for inclusion in a professional report by an engineering consultancy. "
            f"It should be written in way that the paragraph text refers to the figures to illustrate the main technical points. "
            f"Each slide contains technical data, and your task is to highlight key points and any notable differences or observations.\n\n"
            f"Titles: {combined_titles}\n\nText: {combined_text}\n\n"
            f"Generate the description based on this content, ensuring it is suitable for a technical audience."
        )

        if example_report:
            prompt += f"\n\nHere is an example of a report that illustrates the style and quality expected:\n\n{example_report}"

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a technical writer at Three60 Energy."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1500
        )
        return response.choices[0].message['content'].strip()

    # Function to add a caption to a paragraph in a Word document
    def add_caption(paragraph, figure_count, caption_text):
        sanitized_caption_text = sanitize_text(caption_text)
        run = paragraph.add_run()
        fldChar_begin = OxmlElement('w:fldChar')
        fldChar_begin.set(qn('w:fldCharType'), 'begin')

        instrText = OxmlElement('w:instrText')
        instrText.text = f'SEQ Figure \* ARABIC'
        
        fldChar_separate = OxmlElement('w:fldChar')
        fldChar_separate.set(qn('w:fldCharType'), 'separate')

        fldChar_end = OxmlElement('w:fldChar')
        fldChar_end.set(qn('w:fldCharType'), 'end')

        run._r.append(fldChar_begin)
        run._r.append(instrText)
        run._r.append(fldChar_separate)
        run.add_text(f' Figure {figure_count}: {sanitized_caption_text}')
        run._r.append(fldChar_end)

    # Function to parse slide ranges
    def parse_slide_ranges(slide_ranges):
        slides = set()
        ranges = slide_ranges.split(',')
        for r in ranges:
            r = r.strip()
            if '-' in r:
                start, end = map(int, r.split('-'))
                slides.update(range(start, end + 1))
            elif r.isdigit():  # Ensure r is a valid digit
                slides.add(int(r))
        return sorted(slides)

    def create_word_report(report, template_path=None):
        doc = Document(template_path) if template_path else Document()
        figure_count = 1
        first_section = True
        
        for section, content in report.items():
            if section.startswith("Heading 1: "):
                if not first_section:
                    doc.add_page_break()  # Add a page break before each new section
                first_section = False
                heading_text = section.replace("Heading 1: ", "").strip()
                doc.add_heading(heading_text, level=1)
            else:
                heading_text = section.replace("Heading 2: ", "").strip()
                doc.add_heading(heading_text, level=2)
            
            # Iterate over the content to add text and figures sequentially
            for slide in content["slides"]:
                sanitized_text = sanitize_text(slide["text"])
                # Ensure there are no unwanted newlines or spaces
                sanitized_text = " ".join(sanitized_text.split())
                doc.add_paragraph(sanitized_text)
                
                if os.path.exists(slide["screenshot"]):
                    doc.add_picture(slide["screenshot"], width=Inches(4.0))
                    p = doc.add_paragraph()
                    add_caption(p, figure_count, sanitize_text(slide["title"]).strip())
                    figure_count += 1
                    # Add a paragraph with two new lines for spacing
                    doc.add_paragraph().add_run("\n\n")

        return doc

    # Initialize session state for slides and slide_index if not already initialized
    if 'slides' not in st.session_state:
        st.session_state.slides = []

    if 'slide_index' not in st.session_state:
        st.session_state.slide_index = 0

    # Streamlit app
    st.title("PowerPoint to Report Converter")

    uploaded_pptx = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
    uploaded_pdf = st.file_uploader("Upload the corresponding PDF file", type=["pdf"])  # Moved this line here

    presentation_context = st.text_input("Context for the entire presentation")

    uploaded_template = st.file_uploader("Upload a Word template file (optional)", type=["docx"])
    uploaded_report_example = st.file_uploader("Upload an example report (optional)", type=["txt", "docx"])

    example_report = None
    if uploaded_report_example:
        if uploaded_report_example.type == "txt":
            example_report = uploaded_report_example.getvalue().decode("utf-8")
        elif uploaded_report_example.type == "docx":
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp_report_example:
                tmp_report_example.write(uploaded_report_example.getbuffer())
                tmp_report_example_path = tmp_report_example.name
            doc = Document(tmp_report_example_path)
            example_report = "\n".join([para.text for para in doc.paragraphs])

    if uploaded_pptx and uploaded_pdf:  # Ensure both files are uploaded
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(uploaded_pptx.getbuffer())
            tmp_pptx_path = tmp.name

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
            tmp_pdf.write(uploaded_pdf.getbuffer())
            pdf_path = tmp_pdf.name

        if uploaded_template:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp_template:
                tmp_template.write(uploaded_template.getbuffer())
                template_path = tmp_template.name
        else:
            template_path = None

        if 'slides' not in st.session_state or not st.session_state.slides:
            ensure_image_directory()

            # Extract text from PPTX
            slides_text = extract_text_from_pptx(tmp_pptx_path)

            # Define the crop box (left, upper, right, lower) based on your needs
            crop_box = (0, 250, 2500, 1430)  # Adjust these values according to your needs
            slides_images = extract_images_from_pdf(pdf_path, crop_box)

            slides = []
            for i, slide in enumerate(slides_text):
                slide["screenshot"] = slides_images[i]
                slides.append(slide)

            st.session_state.slides = slides
            st.session_state.slide_index = 0

    slides = st.session_state.slides
    if len(slides) > 0:
        st.write(f"Number of extracted slides: {len(slides)}")

        # Sidebar for slide navigation and section assignment
        st.sidebar.header("Slide Navigation")
        if st.sidebar.button("Previous Slide"):
            st.session_state.slide_index = max(st.session_state.slide_index - 1, 0)
        if st.sidebar.button("Next Slide"):
            st.session_state.slide_index = min(st.session_state.slide_index + 1, len(slides) - 1)

        current_slide = slides[st.session_state.slide_index]
        st.sidebar.write(f"Slide {st.session_state.slide_index + 1}")
        st.sidebar.write(current_slide["title"])
        st.sidebar.write(current_slide["text"])
        if os.path.exists(current_slide["screenshot"]):
            st.sidebar.image(current_slide["screenshot"], caption=f"Screenshot from Slide {st.session_state.slide_index + 1}")

        st.write("Classify Sections")

        segment_data = {"Heading 1": ["" for _ in range(10)], "Heading 2": ["" for _ in range(10)], "Slides": ["" for _ in range(10)]}
        segment_df = pd.DataFrame(segment_data)
        edited_df = st.data_editor(segment_df, num_rows="fixed", use_container_width=True)

        if st.button("Generate Report"):
            report = {}
            for idx, row in edited_df.dropna(subset=["Slides"]).iterrows():
                heading1 = row["Heading 1"]
                heading2 = row["Heading 2"]
                slide_ranges = row["Slides"]
                slide_indices = parse_slide_ranges(slide_ranges)
                slide_content = [slides[i - 1] for i in slide_indices]

                if heading1:
                    section_title = f"Heading 1: {heading1}"
                else:
                    section_title = f"Heading 2: {heading2}"

                section_texts = []
                description = generate_section_text(slide_content, section_title, presentation_context)
                section_texts.append(description)
                for slide in slide_content:
                    slide["text"] = description  # Update the slide text with the generated description

                report[section_title] = {"text": "\n\n".join(section_texts), "slides": slide_content}

            doc = create_word_report(report, template_path)
            doc.save("report.docx")
            st.write("Report generated. Download the report below:")
            with open("report.docx", "rb") as f:
                st.download_button("Download Report", data=f, file_name="report.docx")
    else:
        st.write("No slides extracted. Please upload a PowerPoint file and the corresponding PDF file.")
else:
    if st.button("Clear"):
        st.session_state.clear()
